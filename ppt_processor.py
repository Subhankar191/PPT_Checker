from pptx import Presentation
from PIL import Image
from io import BytesIO
import os
from image_processor import ImageProcessor

class PPTProcessor:
    def __init__(self):
        self.image_processor = ImageProcessor()
    
    def process(self, pptx_path, verbose=False):
        presentation = Presentation(pptx_path)
        slides_data = []
        
        for i, slide in enumerate(presentation.slides, 1):
            slide_data = {
                'slide_number': i,
                'text': self._extract_slide_text(slide),
                'images_text': [],
                'notes': self._extract_notes_text(slide)
            }
            
            # Process images in slide
            if verbose:
                print(f"Processing slide {i}...")
            
            slide_data['images_text'] = self._process_slide_images(slide, verbose)
            
            slides_data.append(slide_data)
        
        return slides_data
    
    def _extract_slide_text(self, slide):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text.strip())
        return "\n".join(text)
    
    def _extract_notes_text(self, slide):
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, "notes_text_frame"):
                return notes_slide.notes_text_frame.text.strip()
        return ""
    
    def _process_slide_images(self, slide, verbose=False):
        images_text = []
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture type
                image_stream = BytesIO()
                shape.image.blob.seek(0)
                image_stream.write(shape.image.blob.read())
                image_stream.seek(0)
                
                try:
                    img = Image.open(image_stream)
                    text = self.image_processor.extract_text_from_image(img)
                    if text.strip():
                        images_text.append(text)
                        if verbose:
                            print(f"  Extracted text from image: {text[:50]}...")
                except Exception as e:
                    if verbose:
                        print(f"  Error processing image: {str(e)}")
        
        return images_text